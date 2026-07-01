"""Generates the Phase 1+2 midterm evaluation presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
REPORT = os.path.join(ROOT, "docs", "report")
P1 = os.path.join(REPORT, "phase 1")
P2 = os.path.join(REPORT, "phase 2")

DARK = RGBColor(0x1F, 0x4E, 0x79)
MED = RGBColor(0x2E, 0x74, 0xB5)
LIGHT = RGBColor(0xD6, 0xE4, 0xF0)
GREEN = RGBColor(0x70, 0xAD, 0x47)
AMBER = RGBColor(0xFF, 0xC0, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARKTXT = RGBColor(0x33, 0x33, 0x33)
ALTROW = RGBColor(0xEB, 0xF3, 0xFB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_header_bar(slide, title, subtitle=None, kicker=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    bar.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), SW - Inches(1.0), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    if kicker:
        p0 = tf.paragraphs[0]
        p0.text = kicker
        p0.font.size = Pt(13)
        p0.font.color.rgb = LIGHT
        p0.font.bold = True
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = LIGHT
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), SW, Pt(4))
    accent.fill.solid()
    accent.fill.fore_color.rgb = MED
    accent.line.fill.background()
    accent.shadow.inherit = False
    return bar


def add_footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), SH - Inches(0.4), SW - Inches(1.5), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    pg = tb.text_frame.add_paragraph()


def bullets(slide, items, left, top, width, height, size=18, color=DARKTXT, bold_first=False,
            bullet_color=MED, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        marker = "▸ " if level == 0 else "–  "
        p.text = marker + text
        p.font.size = Pt(size - level * 2)
        p.font.color.rgb = color
        p.line_spacing = line_spacing
        p.space_after = Pt(8 if level == 0 else 4)
    return tb


def title_slide():
    s = add_slide()
    add_bg(s, DARK)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.6), SW, Inches(0.06))
    band.fill.solid(); band.fill.fore_color.rgb = MED; band.line.fill.background(); band.shadow.inherit = False

    tb = s.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = "M.Tech. (Online) — Artificial Intelligence"
    p.font.size = Pt(18); p.font.color.rgb = LIGHT

    tb2 = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.2))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Detecting AI-Generated and Cloned Speech"
    p2.font.size = Pt(40); p2.font.bold = True; p2.font.color.rgb = WHITE
    p3 = tf2.add_paragraph()
    p3.text = "to Strengthen Voice-Based Authentication in Banking Environments"
    p3.font.size = Pt(24); p3.font.color.rgb = LIGHT

    tb3 = s.shapes.add_textbox(Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.6))
    tf3 = tb3.text_frame
    lines = [
        ("Midterm Evaluation & Demo — Phase 1 + Phase 2 Progress", 16, WHITE, True),
        ("Rajib Roy  |  Student No. 24459  |  IISc M.Tech (AI)", 15, LIGHT, False),
        ("Internal Guide: Dr. Anil Rahate, Wipro Ltd   |   IISc Mentor: Prof. C. S. Seelamantula", 13, LIGHT, False),
        ("July 2026", 13, LIGHT, False),
    ]
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = text
        p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
        p.space_after = Pt(6)
    return s


def agenda_slide():
    s = add_slide(); add_bg(s)
    add_header_bar(s, "Agenda", kicker="OVERVIEW")
    items = [
        "Problem & motivation",
        "Research questions and novelty",
        "Project timeline: Phase 1 → Phase 2 → Phase 3",
        "Phase 1 recap: dataset & feature pipeline",
        "Phase 2: classical ML baselines",
        "Phase 2: feature importance analysis",
        "Phase 2: deep learning models (CNN, CNN-LSTM)",
        "Phase 2: robustness & cross-system generalisation",
        "Consolidated results & key insights",
        "Phase 3 plan & Q&A",
    ]
    bullets(s, items, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), size=20)
    add_footer(s, "Midterm Evaluation — Phase 1 + Phase 2")
    return s


def problem_slide():
    s = add_slide(); add_bg(s)
    add_header_bar(s, "The Problem", kicker="MOTIVATION")
    left_items = [
        "Voice authentication is routine in banking IVR / call-centre workflows",
        "Neural TTS & voice cloning can now replicate a voice from a few seconds of audio",
        "Existing voice biometrics were designed against replay attacks — not modern neural TTS",
        "A successful spoof enables fraudulent transfers, identity fraud, social engineering",
    ]
    bullets(s, left_items, Inches(0.8), Inches(1.5), Inches(6.0), Inches(5.0), size=18)

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(1.6), Inches(5.4), Inches(4.6))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.color.rgb = MED; box.line.width = Pt(1)
    tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]; p.text = "Goal"; p.font.bold = True; p.font.size = Pt(20); p.font.color.rgb = DARK
    p2 = tf.add_paragraph()
    p2.text = ("Build and validate an automated pipeline that distinguishes genuine human speech "
               "from AI-generated / cloned speech, specifically under banking telephony conditions: "
               "short utterances, codec compression, background noise.")
    p2.font.size = Pt(16); p2.font.color.rgb = DARKTXT; p2.space_before = Pt(10)
    add_footer(s, "Midterm Evaluation — Phase 1 + Phase 2")
    return s


def rq_slide():
    s = add_slide(); add_bg(s)
    add_header_bar(s, "Research Questions", kicker="SCOPE")
    rqs = [
        ("RQ1", "Can supervised ML reliably distinguish real vs. AI-synthesised speech, on unseen speakers and unseen TTS systems?"),
        ("RQ2", "Which acoustic/signal-level properties (spectral shape, phase coherence, prosody) most reliably betray synthetic speech?"),
        ("RQ3", "How do detection models hold up under banking telephony conditions — short utterances, G.711 codec, ambient noise?"),
        ("RQ4", "Do deep learning architectures offer a meaningful advantage over classical ML under identical data conditions?"),
    ]
    top = Inches(1.55)
    for tag, text in rqs:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.15))
        card.fill.solid(); card.fill.fore_color.rgb = ALTROW; card.line.color.rgb = MED; card.line.width = Pt(0.75)
        tag_box = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), top + Inches(0.2), Inches(0.75), Inches(0.75))
        tag_box.fill.solid(); tag_box.fill.fore_color.rgb = DARK; tag_box.line.fill.background()
        tf0 = tag_box.text_frame; tf0.paragraphs[0].text = tag
        tf0.paragraphs[0].font.size = Pt(14); tf0.paragraphs[0].font.bold = True
        tf0.paragraphs[0].font.color.rgb = WHITE; tf0.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb = s.shapes.add_textbox(Inches(2.0), top + Inches(0.12), Inches(10.3), Inches(0.95))
        tfb = tb.text_frame; tfb.word_wrap = True
        pb = tfb.paragraphs[0]; pb.text = text; pb.font.size = Pt(15); pb.font.color.rgb = DARKTXT
        top += Inches(1.32)
    add_footer(s, "Midterm Evaluation — Phase 1 + Phase 2")
    return s


def timeline_slide():
    s = add_slide(); add_bg(s)
    add_header_bar(s, "Project Timeline", kicker="WORK PLAN (APPROVED PROPOSAL)")
    phases = [
        ("Phase 1", "Jan – Apr 2026", "Literature review · Dataset (17,947 utterances) · Feature pipeline", GREEN, "COMPLETE"),
        ("Phase 2", "May – Jul 2026", "ML baselines · Feature importance · CNN / CNN-LSTM · Robustness experiments", GREEN, "COMPLETE"),
        ("Phase 3", "Aug – Nov 2026", "SHAP + saliency interpretability · Final evaluation & stats · Thesis writing", AMBER, "UPCOMING"),
    ]
    left = Inches(0.8)
    w = Inches(3.85)
    top = Inches(1.9)
    for i, (name, dates, desc, color, status) in enumerate(phases):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + i * (w + Inches(0.25)), top, w, Inches(4.5))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = DARK; card.line.width = Pt(1.25)
        tf = card.text_frame; tf.word_wrap = True; tf.margin_top = Inches(0.25); tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.2)
        p1 = tf.paragraphs[0]; p1.text = name; p1.font.size = Pt(24); p1.font.bold = True; p1.font.color.rgb = DARK
        p2 = tf.add_paragraph(); p2.text = dates; p2.font.size = Pt(14); p2.font.color.rgb = MED; p2.space_after = Pt(14)
        p3 = tf.add_paragraph(); p3.text = desc; p3.font.size = Pt(14); p3.font.color.rgb = DARKTXT; p3.line_spacing = 1.25
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + i * (w + Inches(0.25)) + Inches(0.25),
                                    top + Inches(3.85), Inches(1.7), Inches(0.4))
        badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
        bt = badge.text_frame.paragraphs[0]; bt.text = status; bt.font.size = Pt(12); bt.font.bold = True
        bt.font.color.rgb = WHITE; bt.alignment = PP_ALIGN.CENTER
    add_footer(s, "We are here: end of Phase 2, entering Phase 3")
    return s


def dataset_slide():
    s = add_slide(); add_bg(s)
    add_header_bar(s, "Phase 1 Recap: Dataset", kicker="FOUNDATION FOR PHASE 2")
    left_items = [
        "17,947 utterances from 10 sources",
        ("Genuine: LibriSpeech test-clean, 2,229 clips, 40 speakers", 1),
        ("Synthetic: ASVspoof 2019 LA (6 TTS/VC systems), VITS, Tacotron2, Edge-TTS", 1),
        "8 kHz + G.711 codec simulation (banking telephony)",
        "3–8 second utterances (transaction-confirmation length)",
        "Strict speaker- and system-independent splits",
        "Microsoft Edge-TTS held out entirely — unseen-system generalisation test",
    ]
    bullets(s, left_items, Inches(0.8), Inches(1.55), Inches(6.3), Inches(5.2), size=17)
    img = os.path.join(P1, "Fig2_Dataset_Composition.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(7.2), Inches(1.7), width=Inches(5.4))
    add_footer(s, "Midterm Evaluation — Phase 1 + Phase 2")
    return s


def features_slide():
    s = add_slide(); add_bg(s)
    add_header_bar(s, "Phase 1 Recap: Feature Pipeline", kicker="FOUNDATION FOR PHASE 2")
    left_items = [
        "257-dim handcrafted feature vector (Librosa, 25 ms window / 10 ms hop)",
        ("MFCC (40 coeff. × static/Δ/ΔΔ) = 120 dims", 1),
        ("Spectral centroid, bandwidth, rolloff, ZCR, HNR, pitch, phase difference", 1),
        "128-bin log-Mel spectrograms for deep learning models (128×128)",
        "All hyperparameters centralised in config.yaml",
        "diagnose.py validates class balance, NaNs, inter-class separation",
    ]
    bullets(s, left_items, Inches(0.8), Inches(1.55), Inches(6.3), Inches(5.2), size=17)
    img = os.path.join(P1, "Fig3_Feature_Vector.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(7.2), Inches(1.7), width=Inches(5.4))
    add_footer(s, "Midterm Evaluation — Phase 1 + Phase 2")
    return s


def table_slide(title, kicker, headers, rows, col_widths, highlight_rows=None, note=None, footer=""):
    s = add_slide(); add_bg(s)
    add_header_bar(s, title, kicker=kicker)
    highlight_rows = highlight_rows or []
    n_rows = len(rows) + 1
    n_cols = len(headers)
    left = Inches(0.8)
    top = Inches(1.55)
    total_w = Inches(11.7)
    tbl_height = Inches(0.5) * n_rows
    graphic_frame = s.shapes.add_table(n_rows, n_cols, left, top, total_w, tbl_height)
    tbl = graphic_frame.table
    for c, w in enumerate(col_widths):
        tbl.columns[c].width = w
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = DARK
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREEN if r in highlight_rows else (ALTROW if r % 2 == 0 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12.5)
            p.font.bold = r in highlight_rows
            p.font.color.rgb = WHITE if r in highlight_rows else DARKTXT
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if note:
        tb = s.shapes.add_textbox(left, top + tbl_height + Inches(0.15), total_w, Inches(0.9))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = note; p.font.size = Pt(14); p.font.color.rgb = DARKTXT; p.font.italic = True
    add_footer(s, footer or "Midterm Evaluation — Phase 1 + Phase 2")
    return s


def baseline_results_slide():
    headers = ["Model", "Eval. Set", "EER ↓", "F1 ↑", "AUC ↑", "Acc. ↑"]
    rows = [
        ["SVM (RBF)", "Test", "0.0083", "0.9943", "0.9992", "0.9912"],
        ["SVM (RBF)", "Gen. test†", "0.0067", "0.9969", "0.9987", "0.9953"],
        ["Random Forest", "Test", "0.0305", "0.9834", "0.9956", "0.9744"],
        ["Random Forest", "Gen. test†", "0.0293", "0.9854", "0.9964", "0.9781"],
        ["Gradient Boosting", "Test", "0.0346", "0.9820", "0.9949", "0.9722"],
        ["Gradient Boosting", "Gen. test†", "0.0311", "0.9868", "0.9956", "0.9801"],
    ]
    cw = [Inches(2.6), Inches(2.0), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7)]
    note = ("SVM (RBF kernel, C=10) is the strongest classical baseline. All three models trained with "
            "SMOTE + class-weighting inside 5-fold CV. † Gen. test = held-out Edge-TTS, unseen in training.")
    return table_slide("Phase 2: Classical ML Baselines", "SVM · RANDOM FOREST · GRADIENT BOOSTING",
                        headers, rows, cw, highlight_rows=[0], note=note)


def feature_importance_slide():
    s = add_slide(); add_bg(s)
    add_header_bar(s, "Phase 2: Feature Importance Analysis", kicker="INTERPRETABILITY (RANDOM FOREST)")
    items = [
        "Random Forest importance ranking over the 257-dim feature vector",
        "Delta-MFCC coefficients dominate (delta_mfcc_6_mean, delta_mfcc_5_mean)",
        ("→ TTS vocoders impose smoother, more regular spectral dynamics than natural articulation", 1),
        "phase_diff_mean ranks 16th — validates inclusion of phase-based features",
        "Spectral centroid & rolloff appear as secondary discriminators",
    ]
    bullets(s, items, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8), size=17)
    img = os.path.join(P2, "feature_importance.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(6.7), Inches(1.55), width=Inches(6.0))
    add_footer(s, "Midterm Evaluation — Phase 1 + Phase 2")
    return s


def dl_architecture_slide():
    s = add_slide(); add_bg(s)
    add_header_bar(s, "Phase 2: Deep Learning Models", kicker="CNN & CNN-LSTM ON LOG-MEL SPECTROGRAMS")

    def arch_card(x, name, detail_lines, params):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.55), Inches(5.7), Inches(4.9))
        card.fill.solid(); card.fill.fore_color.rgb = ALTROW; card.line.color.rgb = MED; card.line.width = Pt(1)
        tf = card.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.25); tf.margin_right = Inches(0.25)
        p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = DARK
        for line in detail_lines:
            pp = tf.add_paragraph(); pp.text = "▸ " + line; pp.font.size = Pt(14.5); pp.font.color.rgb = DARKTXT
            pp.space_after = Pt(8); pp.line_spacing = 1.15
        pparam = tf.add_paragraph(); pparam.text = params; pparam.font.size = Pt(15); pparam.font.bold = True
        pparam.font.color.rgb = MED; pparam.space_before = Pt(10)
        return card

    arch_card(Inches(0.8), "CNN",
              ["4 conv blocks: 32→64→128→256 channels, 3×3 kernel, BatchNorm, ReLU, max-pool",
               "Global average pooling",
               "Classifier: 256→128→2, dropout p=0.5"],
              "421,954 parameters")
    arch_card(Inches(6.85), "CNN-LSTM",
              ["Same conv front-end, frequency-only pooling (preserves time axis)",
               "2-layer bidirectional LSTM, hidden size 128/direction",
               "Same classifier head — explicitly models temporal prosodic dynamics"],
              "1,834,498 parameters")

    tb = s.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = ("Training: Adam (lr=1e-3), cosine annealing, 50 epochs, weighted cross-entropy, "
              "WeightedRandomSampler, mixed precision on RTX 3060 (6.4 GB), early stopping (patience 10)")
    p.font.size = Pt(13); p.font.italic = True; p.font.color.rgb = DARKTXT
    return s


def dl_results_slide():
    headers = ["Model", "Params", "Best Epoch", "Val EER", "Training Time"]
    rows = [
        ["CNN", "421,954", "22", "0.0154", "~55 min (RTX 3060)"],
        ["CNN-LSTM", "1,834,498", "21", "0.0092", "~60 min (RTX 3060)"],
    ]
    cw = [Inches(2.5), Inches(2.3), Inches(2.3), Inches(2.3), Inches(2.3)]
    note = ("CNN-LSTM converges faster to a lower validation EER — explicit temporal modelling via "
            "the LSTM captures prosodic artefacts that CNN global average pooling discards.")
    s = table_slide("Phase 2: Deep Learning — Convergence", "CNN VS. CNN-LSTM", headers, rows, cw,
                     highlight_rows=[1], note=note)
    img = os.path.join(P2, "roc_curves.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(1.6), Inches(3.4), width=Inches(10.1))
    return s


def confusion_slide():
    s = add_slide(); add_bg(s)
    add_header_bar(s, "Phase 2: Confusion Matrices", kicker="STANDARD TEST SET")
    img = os.path.join(P2, "confusion_matrices.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(1.2), Inches(1.5), width=Inches(10.9))
    tb = s.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "CNN and CNN-LSTM produce fewer than 3 false positives and under 50 false negatives per 2,734 test clips."
    p.font.size = Pt(14); p.font.italic = True; p.font.color.rgb = DARKTXT
    return s


def robustness_slide():
    s = add_slide(); add_bg(s)
    add_header_bar(s, "Phase 2: Robustness Evaluation", kicker="NOISE (SNR) & G.711 CODEC COMPRESSION")
    img = os.path.join(P2, "degradation_curves.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(1.0), Inches(1.5), width=Inches(11.3))
    items_text = ("SVM degrades fastest under noise (EER ≈0.45 at 5 dB SNR)  ·  Gradient Boosting is most "
                  "noise-robust (flattens at ≈0.33 below 15 dB)  ·  G.711 alone degrades all models only to "
                  "EER 0.05–0.08 → noise is the bigger deployment risk, not codec compression.")
    tb = s.shapes.add_textbox(Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = items_text; p.font.size = Pt(13.5); p.font.color.rgb = DARKTXT
    return s


def consolidated_results_slide():
    headers = ["Model", "Eval. Set", "EER ↓", "F1 ↑", "AUC ↑", "Acc. ↑"]
    rows = [
        ["SVM", "Test", "0.0083", "0.9943", "0.9992", "0.9912"],
        ["SVM", "Gen. test†", "0.0067", "0.9969", "0.9987", "0.9953"],
        ["Random Forest", "Test", "0.0305", "0.9834", "0.9956", "0.9744"],
        ["Random Forest", "Gen. test†", "0.0293", "0.9854", "0.9964", "0.9781"],
        ["Gradient Boosting", "Test", "0.0346", "0.9820", "0.9949", "0.9722"],
        ["Gradient Boosting", "Gen. test†", "0.0311", "0.9868", "0.9956", "0.9801"],
        ["CNN", "Test", "0.0081", "0.9881", "0.9993", "0.9817"],
        ["CNN", "Gen. test†", "0.0036", "0.9942", "0.9991", "0.9914"],
        ["CNN-LSTM", "Test", "0.0035", "0.9965", "0.9996", "0.9945"],
        ["CNN-LSTM", "Gen. test†", "0.0004", "0.9978", "1.0000", "0.9967"],
    ]
    cw = [Inches(2.4), Inches(2.0), Inches(1.75), Inches(1.75), Inches(1.75), Inches(1.75)]
    note = ("CNN-LSTM is best overall (EER 0.0035 test, 0.0004 on unseen Edge-TTS). Every model generalises "
            "as well as or better than on the standard test set. † Gen. test = held-out Edge-TTS.")
    return table_slide("Consolidated Results — All Models", "FULL COMPARISON ACROSS EVALUATION CONDITIONS",
                        headers, rows, cw, highlight_rows=[8, 9], note=note)


def insights_slide():
    s = add_slide(); add_bg(s)
    add_header_bar(s, "Key Insights", kicker="WHAT PHASE 2 TAUGHT US")
    cards = [
        ("Deep Learning vs. Classical", "CNN-LSTM wins overall, but SVM (EER 0.0083) is surprisingly "
         "competitive — even beating the plain CNN. The value of deep learning here comes specifically "
         "from temporal (LSTM) modelling, not convolution alone."),
        ("Generalisation is stronger than expected", "Every model does as well or better on the fully "
         "unseen Edge-TTS system as on the standard test set — spectral/temporal artefacts appear to be "
         "common across neural TTS architectures, not system-specific fingerprints."),
        ("Noise > Codec Compression as a risk", "G.711 alone only pushes EER to 0.05–0.08. Additive noise "
         "pushes classical model EER above 0.3 at low SNR — noise-aware training is the priority for "
         "Phase 3 hardening, ahead of further codec-specific work."),
    ]
    top = Inches(1.6)
    for name, text in cards:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.6))
        card.fill.solid(); card.fill.fore_color.rgb = ALTROW; card.line.color.rgb = MED; card.line.width = Pt(0.75)
        tf = card.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(17); p.font.bold = True; p.font.color.rgb = DARK
        p2 = tf.add_paragraph(); p2.text = text; p2.font.size = Pt(14); p2.font.color.rgb = DARKTXT
        p2.line_spacing = 1.15; p2.space_before = Pt(4)
        top += Inches(1.75)
    add_footer(s, "Midterm Evaluation — Phase 1 + Phase 2")
    return s


def phase3_slide():
    s = add_slide(); add_bg(s)
    add_header_bar(s, "Phase 3 Plan", kicker="AUGUST – NOVEMBER 2026")
    headers = ["Task", "Timeline", "Key Deliverable"]
    rows = [
        ["Interpretability: SHAP for classical models", "Aug 2026", "SHAP summary plots, cross-checked vs. RF importances"],
        ["Interpretability: gradient saliency (CNN/CNN-LSTM)", "Aug–Sep 2026", "Time-frequency saliency maps"],
        ["Final evaluation & statistical testing", "Sep–Oct 2026", "Significance testing between models"],
        ["Thesis writing & submission", "Oct–Nov 2026", "Final report, institute-prescribed format"],
    ]
    cw = [Inches(5.2), Inches(2.2), Inches(4.3)]
    left = Inches(0.8); top = Inches(1.6)
    graphic_frame = s.shapes.add_table(len(rows) + 1, 3, left, top, Inches(11.7), Inches(0.55) * (len(rows) + 1))
    tbl = graphic_frame.table
    for c, w in enumerate(cw):
        tbl.columns[c].width = w
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = DARK
        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER; cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c); cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = ALTROW if r % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]; p.font.size = Pt(13); p.font.color.rgb = DARKTXT
            p.alignment = PP_ALIGN.LEFT; cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_footer(s, "Midterm Evaluation — Phase 1 + Phase 2")
    return s


def closing_slide():
    s = add_slide(); add_bg(s, DARK)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Thank you"
    p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE
    p2 = tb.text_frame.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.font.size = Pt(22); p2.font.color.rgb = LIGHT
    return s


title_slide()
agenda_slide()
problem_slide()
rq_slide()
timeline_slide()
dataset_slide()
features_slide()
baseline_results_slide()
feature_importance_slide()
dl_architecture_slide()
dl_results_slide()
confusion_slide()
robustness_slide()
consolidated_results_slide()
insights_slide()
phase3_slide()
closing_slide()

out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Phase2_Midterm_Presentation.pptx")
prs.save(out_path)
print("Saved:", out_path)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
